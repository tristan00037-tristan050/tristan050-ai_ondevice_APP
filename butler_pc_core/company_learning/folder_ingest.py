from __future__ import annotations

import email
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from butler_pc_core.company_policy.contracts import sha256_bytes, sha256_text, stable_json_digest
from butler_pc_core.retrieval.chunkers import Chunk, get_chunker
from butler_pc_core.retrieval.pipeline import PersonalPackIndex

ALLOWED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".csv",
    ".pptx",
    ".ppt",
    ".eml",
    ".msg",
}
DEFAULT_MAX_DEPTH = 6
DEFAULT_MAX_FILES = 10_000
DEFAULT_MAX_FILE_BYTES = 25 * 1024 * 1024
DEFAULT_MAX_TOTAL_BYTES = 1024 * 1024 * 1024
DEFAULT_MAX_EXTRACTED_CHARS_PER_FILE = 4_000_000
DEFAULT_MAX_TOTAL_EXTRACTED_CHARS = 64_000_000
MAX_PDF_PAGES = 2_000
MAX_OOXML_MEMBERS = 10_000
MAX_OOXML_UNCOMPRESSED_BYTES = 128 * 1024 * 1024
MAX_OOXML_COMPRESSION_RATIO = 250.0
MAX_ARCHIVE_ENTRIES = 10_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 256 * 1024 * 1024
MAX_ARCHIVE_MEMBER_BYTES = 128 * 1024 * 1024

_SKIP_DIR_NAMES = {
    "__pycache__",
    ".cache",
    ".git",
    ".hg",
    ".svn",
    "node_modules",
    "tmp",
    "temp",
    "cache",
    "Library",
}
_SKIP_SUFFIXES = (".tmp", ".temp", ".lock", ".part", ".crdownload")


class CompanyLearningIngestError(ValueError):
    def __init__(self, fail_class: str) -> None:
        super().__init__(fail_class)
        self.fail_class = fail_class


@dataclass(frozen=True)
class IngestCounters:
    scanned_files: int = 0
    processed_files: int = 0
    skipped_unsupported: int = 0
    skipped_hidden: int = 0
    skipped_symlink: int = 0
    skipped_limit: int = 0
    skipped_extractor: int = 0
    failed_extract: int = 0
    total_bytes: int = 0
    chunk_count: int = 0
    by_extension: dict[str, int] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "scanned_files": self.scanned_files,
            "processed_files": self.processed_files,
            "skipped_unsupported": self.skipped_unsupported,
            "skipped_hidden": self.skipped_hidden,
            "skipped_symlink": self.skipped_symlink,
            "skipped_limit": self.skipped_limit,
            "skipped_extractor": self.skipped_extractor,
            "failed_extract": self.failed_extract,
            "total_bytes": self.total_bytes,
            "chunk_count": self.chunk_count,
            "by_extension": dict(sorted(self.by_extension.items())),
        }


@dataclass(frozen=True)
class IngestResult:
    folder_digest: str
    ingest_digest: str
    counters: IngestCounters
    chunks: list[Chunk]
    index: PersonalPackIndex


def _mutate_counter(counters: IngestCounters, **updates: int) -> IngestCounters:
    data = counters.to_dict()
    by_extension = dict(data.pop("by_extension"))
    data.update({key: data[key] + value for key, value in updates.items()})
    return IngestCounters(**data, by_extension=by_extension)


def _with_extension(counters: IngestCounters, ext: str) -> IngestCounters:
    data = counters.to_dict()
    by_extension = dict(data.pop("by_extension"))
    by_extension[ext] = by_extension.get(ext, 0) + 1
    return IngestCounters(**data, by_extension=by_extension)


def _decode_bytes(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _validate_signature(data: bytes, suffix: str) -> None:
    if suffix.lower() in {".docx", ".xlsx", ".pptx"} and not data.startswith(
        (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")
    ):
        raise CompanyLearningIngestError("ARCHIVE_SIGNATURE_INVALID")


def _guard_ooxml_archive(data: bytes) -> None:
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            members = archive.infolist()
    except zipfile.BadZipFile as exc:
        raise CompanyLearningIngestError("ARCHIVE_INVALID") from exc

    if len(members) > min(MAX_OOXML_MEMBERS, MAX_ARCHIVE_ENTRIES):
        raise CompanyLearningIngestError("ARCHIVE_LIMIT_EXCEEDED")

    total_uncompressed = 0
    total_compressed = 0
    for member in members:
        total_uncompressed += int(member.file_size or 0)
        total_compressed += int(member.compress_size or 0)
        if int(member.file_size or 0) > MAX_ARCHIVE_MEMBER_BYTES:
            raise CompanyLearningIngestError("ARCHIVE_LIMIT_EXCEEDED")
        if total_uncompressed > min(MAX_OOXML_UNCOMPRESSED_BYTES, MAX_ARCHIVE_UNCOMPRESSED_BYTES):
            raise CompanyLearningIngestError("ARCHIVE_LIMIT_EXCEEDED")

    if total_uncompressed <= 0:
        return
    if total_compressed <= 0:
        raise CompanyLearningIngestError("ARCHIVE_COMPRESSION_RATIO_EXCEEDED")
    if total_uncompressed / total_compressed > MAX_OOXML_COMPRESSION_RATIO:
        raise CompanyLearningIngestError("ARCHIVE_COMPRESSION_RATIO_EXCEEDED")


def _guarded_ooxml_bytes(path: Path) -> bytes:
    data = path.read_bytes()
    _validate_signature(data, path.suffix)
    _guard_ooxml_archive(data)
    return data


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - depends on runtime package set
        raise CompanyLearningIngestError("EXTRACTOR_UNAVAILABLE") from exc
    try:
        doc = Document(BytesIO(_guarded_ooxml_bytes(path)))
        parts: list[str] = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
        return "\n".join(part for part in parts if part.strip())
    except CompanyLearningIngestError:
        raise
    except Exception as exc:
        raise CompanyLearningIngestError("EXTRACT_FAILED") from exc


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - depends on runtime package set
        raise CompanyLearningIngestError("EXTRACTOR_UNAVAILABLE") from exc
    try:
        wb = load_workbook(BytesIO(_guarded_ooxml_bytes(path)), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets[:20]:
            for row in sheet.iter_rows(max_row=500, values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines)
    except CompanyLearningIngestError:
        raise
    except Exception as exc:
        raise CompanyLearningIngestError("EXTRACT_FAILED") from exc


def _extract_pdf(path: Path) -> str:
    try:
        from pdfminer.high_level import extract_text
        from pdfminer.pdfpage import PDFPage
    except ImportError as exc:  # pragma: no cover - depends on runtime package set
        raise CompanyLearningIngestError("EXTRACTOR_UNAVAILABLE") from exc
    try:
        with path.open("rb") as handle:
            for page_index, _page in enumerate(PDFPage.get_pages(handle, caching=False), start=1):
                if page_index > MAX_PDF_PAGES:
                    raise CompanyLearningIngestError("PDF_PAGE_LIMIT_EXCEEDED")
        return str(extract_text(str(path)) or "")
    except CompanyLearningIngestError:
        raise
    except Exception as exc:
        raise CompanyLearningIngestError("EXTRACT_FAILED") from exc


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover - optional local dependency
        raise CompanyLearningIngestError("EXTRACTOR_UNAVAILABLE") from exc
    try:
        prs = Presentation(BytesIO(_guarded_ooxml_bytes(path)))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    parts.append(text.strip())
        return "\n".join(parts)
    except CompanyLearningIngestError:
        raise
    except Exception as exc:
        raise CompanyLearningIngestError("EXTRACT_FAILED") from exc


def _extract_eml(path: Path) -> str:
    try:
        msg = email.message_from_bytes(path.read_bytes())
        parts: list[str] = []
        subject = str(msg.get("subject") or "").strip()
        if subject:
            parts.append(subject)
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() != "text/plain":
                    continue
                payload = part.get_payload(decode=True)
                if isinstance(payload, bytes):
                    parts.append(_decode_bytes(payload))
        else:
            payload = msg.get_payload(decode=True)
            if isinstance(payload, bytes):
                parts.append(_decode_bytes(payload))
        return "\n".join(part for part in parts if part.strip())
    except Exception as exc:
        raise CompanyLearningIngestError("EXTRACT_FAILED") from exc


def extract_runtime_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".csv"}:
        return _decode_bytes(path.read_bytes())
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".eml":
        return _extract_eml(path)
    raise CompanyLearningIngestError("EXTRACTOR_UNAVAILABLE")


def _safe_relative(path: Path, root: Path) -> Path:
    try:
        return path.relative_to(root)
    except ValueError as exc:
        raise CompanyLearningIngestError("FOLDER_ESCAPE_BLOCKED") from exc


def _should_skip_hidden_or_system(path: Path, root: Path) -> bool:
    rel = _safe_relative(path, root)
    if any(part.startswith(".") for part in rel.parts):
        return True
    if any(part in _SKIP_DIR_NAMES for part in rel.parts):
        return True
    return path.name.lower().endswith(_SKIP_SUFFIXES)


def _source_id_for(file_digest: str) -> str:
    return "file:" + file_digest.split(":", 1)[1][:16]


def _chunker_for_allowed_extension(path: Path):
    if path.suffix.lower() in {".txt", ".md"}:
        return get_chunker(file_type="meeting_minutes")
    return get_chunker(file_path=path)


def ingest_folder(
    folder_path: str,
    *,
    max_depth: int = DEFAULT_MAX_DEPTH,
    max_files: int = DEFAULT_MAX_FILES,
    max_file_bytes: int = DEFAULT_MAX_FILE_BYTES,
    max_total_bytes: int = DEFAULT_MAX_TOTAL_BYTES,
    max_extracted_chars_per_file: int = DEFAULT_MAX_EXTRACTED_CHARS_PER_FILE,
    max_total_extracted_chars: int = DEFAULT_MAX_TOTAL_EXTRACTED_CHARS,
) -> IngestResult:
    root = Path(folder_path).expanduser().resolve()
    if not root.exists() or not root.is_dir():
        raise CompanyLearningIngestError("FOLDER_NOT_FOUND")

    folder_digest = sha256_text(str(root))
    counters = IngestCounters()
    chunks: list[Chunk] = []
    file_digests: list[str] = []
    processed_files = 0
    total_bytes = 0
    total_extracted_chars = 0

    for path in sorted(root.rglob("*")):
        if path.is_symlink():
            counters = _mutate_counter(counters, skipped_symlink=1)
            continue
        if not path.is_file():
            continue
        rel = _safe_relative(path.resolve(), root)
        if len(rel.parts) - 1 > max_depth:
            counters = _mutate_counter(counters, skipped_limit=1)
            continue
        if _should_skip_hidden_or_system(path, root):
            counters = _mutate_counter(counters, skipped_hidden=1)
            continue

        counters = _mutate_counter(counters, scanned_files=1)
        ext = path.suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            counters = _mutate_counter(counters, skipped_unsupported=1)
            continue
        counters = _with_extension(counters, ext)

        try:
            size = path.stat().st_size
        except OSError:
            counters = _mutate_counter(counters, failed_extract=1)
            continue
        if size > max_file_bytes or processed_files >= max_files or total_bytes + size > max_total_bytes:
            counters = _mutate_counter(counters, skipped_limit=1)
            continue

        try:
            raw_bytes = path.read_bytes()
            file_digest = sha256_bytes(raw_bytes)
            text = extract_runtime_text(path)
            extracted_chars = len(text)
            if (
                extracted_chars > max_extracted_chars_per_file
                or total_extracted_chars + extracted_chars > max_total_extracted_chars
            ):
                counters = _mutate_counter(counters, skipped_limit=1)
                continue
            if not text.strip():
                counters = _mutate_counter(counters, skipped_extractor=1)
                continue
            chunker = _chunker_for_allowed_extension(path)
            source_id = _source_id_for(file_digest)
            produced = chunker.chunk(text, source_file=source_id)
        except CompanyLearningIngestError as exc:
            if exc.fail_class == "EXTRACTOR_UNAVAILABLE":
                counters = _mutate_counter(counters, skipped_extractor=1)
            else:
                counters = _mutate_counter(counters, failed_extract=1)
            continue
        except Exception:
            counters = _mutate_counter(counters, failed_extract=1)
            continue

        if not produced:
            counters = _mutate_counter(counters, skipped_extractor=1)
            continue

        chunks.extend(produced)
        file_digests.append(file_digest)
        processed_files += 1
        total_bytes += size
        total_extracted_chars += len(text)
        counters = _mutate_counter(
            counters,
            processed_files=1,
            total_bytes=size,
            chunk_count=len(produced),
        )

    if not chunks:
        raise CompanyLearningIngestError("NO_INDEXABLE_FILES")

    ingest_digest = stable_json_digest(
        {
            "folder_digest": folder_digest,
            "file_digests": sorted(file_digests),
            "chunk_digests": [sha256_text(chunk.text) for chunk in chunks],
        }
    )
    index = PersonalPackIndex.build(chunks)
    return IngestResult(
        folder_digest=folder_digest,
        ingest_digest=ingest_digest,
        counters=counters,
        chunks=chunks,
        index=index,
    )
